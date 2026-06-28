#!/usr/bin/env python3
"""Build the 'Beyond Prompting' deck as an editable PowerPoint.

Styled to the BOS brand guide. Uses the brand's declared fallback fonts
(Georgia for serif display, Helvetica Neue for body/labels) so it renders
correctly on any Mac without installing Cormorant Garamond / Jost.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

HERE = os.path.dirname(os.path.abspath(__file__))
ASSETS = os.path.join(HERE, "assets")

# ---- Brand palette ----
CREAM    = RGBColor(0xF5, 0xF0, 0xE8)
SAND     = RGBColor(0xE8, 0xDC, 0xC8)
CLAY     = RGBColor(0xC4, 0x95, 0x6A)
CLAY_LT  = RGBColor(0xF0, 0xE0, 0xCC)
BROWN    = RGBColor(0x3D, 0x2B, 0x1F)
BROWN_MID= RGBColor(0x6B, 0x4C, 0x3B)
COPPER   = RGBColor(0xA0, 0x52, 0x2D)
LIGHT    = RGBColor(0xFA, 0xF7, 0xF2)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)

# Brand fonts -> guaranteed fallbacks (per brand guide font stacks)
SERIF = "Georgia"          # Cormorant Garamond -> Georgia, serif
SANS  = "Helvetica Neue"   # Jost -> Helvetica Neue, sans-serif

EMU_IN = 914400
SW = 13.333
SH = 7.5

prs = Presentation()
prs.slide_width  = Inches(SW)
prs.slide_height = Inches(SH)
BLANK = prs.slide_layouts[6]


def add_slide(bg):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg
    return s


def rect(slide, l, t, w, h, color, shape=MSO_SHAPE.RECTANGLE, line=None, line_w=None):
    sp = slide.shapes.add_shape(shape, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w or 1)
    sp.shadow.inherit = False
    return sp


def set_spc(run, pts):
    """Letter spacing in points (OOXML spc = hundredths of a point)."""
    run._r.get_or_add_rPr().set("spc", str(int(pts * 100)))


def textbox(slide, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    return tb, tf


def para(tf, runs, size, font=SANS, color=BROWN, bold=False, italic=False,
         align=PP_ALIGN.LEFT, space_before=0, space_after=0, line=1.0,
         spc=None, upper=False, first=False):
    """runs: str OR list of (text, overrides-dict)."""
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    p.line_spacing = line
    if isinstance(runs, str):
        runs = [(runs, {})]
    for text, ov in runs:
        r = p.add_run()
        r.text = text.upper() if upper else text
        f = r.font
        f.name = ov.get("font", font)
        f.size = Pt(ov.get("size", size))
        f.bold = ov.get("bold", bold)
        f.italic = ov.get("italic", italic)
        f.color.rgb = ov.get("color", color)
        sp = ov.get("spc", spc)
        if sp is not None:
            set_spc(r, sp)
    return p


def eyebrow(slide, text, l, t, color=CLAY):
    _, tf = textbox(slide, l, t, 9, 0.4)
    para(tf, text, 13, font=SANS, color=color, bold=True, upper=True, spc=2.4, first=True)


def footer(slide, left_txt, right_txt, color=CLAY):
    _, tf = textbox(slide, 0.9, SH - 0.62, 5.5, 0.35)
    para(tf, left_txt, 9.5, font=SANS, color=color, bold=True, upper=True, spc=1.8, first=True)
    _, tf2 = textbox(slide, SW - 6.4, SH - 0.62, 5.5, 0.35)
    para(tf2, right_txt, 9.5, font=SANS, color=color, bold=True, upper=True,
         spc=1.8, align=PP_ALIGN.RIGHT, first=True)


def accent_bar(slide):
    rect(slide, 0, 0, SW, 0.07, CLAY)


def wordmark(slide, l, t, color, scale=1.0):
    """Recreate the BOS wordmark: vertical rule + THE / BUILDERS' / OPS STUDIO."""
    # vertical rule
    rect(slide, l, t + 0.05 * scale, 0.015, 1.35 * scale, color)
    rect(slide, l - 0.08 * scale, t + 0.55 * scale, 0.22 * scale, 0.014, color)
    tx = l + 0.45 * scale
    _, tf = textbox(slide, tx, t, 6, 1.6 * scale)
    para(tf, "THE", 12 * scale, font=SANS, color=color, bold=True, spc=4, first=True)
    para(tf, "BUILDERS'", 26 * scale, font=SERIF, color=color, bold=True, spc=1.2,
         space_before=2, line=0.95)
    para(tf, "OPS STUDIO", 12 * scale, font=SANS, color=color, bold=True, spc=4,
         space_before=2)


# ============================================================
# SLIDE 1 — COVER
# ============================================================
s = add_slide(BROWN)
accent_bar(s)
wordmark(s, 0.95, 1.0, CREAM, scale=1.0)
eyebrow(s, "A Live Build Session", 0.95, 3.05)
_, tf = textbox(s, 0.9, 3.45, 11.5, 2.5)
para(tf, "Beyond", 76, font=SERIF, color=CREAM, bold=True, line=0.98, first=True)
para(tf, "Prompting", 76, font=SERIF, color=CREAM, bold=True, line=0.98)
_, tf = textbox(s, 0.95, 5.85, 9.6, 1.1)
para(tf, [("Building with AI for small business. Not the prompts. The ", {}),
          ("machine", {"color": CLAY, "italic": True, "font": SERIF, "size": 22}),
          (" that runs the back office while you sleep.", {})],
     20, font=SANS, color=SAND, line=1.35, first=True)
footer(s, "Cynthia Jones", "thebuildersopsstudio.com")

# ============================================================
# SLIDE 2 — THE GAP
# ============================================================
s = add_slide(CREAM)
accent_bar(s)
eyebrow(s, "Start Here", 0.95, 1.15)
_, tf = textbox(s, 0.9, 1.6, 11.5, 2.4)
para(tf, [("There's a gap between ", {}), ("using", {"color": COPPER}), (" AI", {})],
     50, font=SERIF, color=BROWN, bold=True, line=1.05, first=True)
para(tf, [("and AI ", {}), ("doing the work.", {"color": COPPER})],
     50, font=SERIF, color=BROWN, bold=True, line=1.05)
_, tf = textbox(s, 0.95, 4.35, 10.6, 1.7)
para(tf, "Most people have used AI to write a caption, an email, a post. Far fewer "
        "have an AI tool that actually ran a part of their business for them, while "
        "they slept.", 20, font=SANS, color=BROWN_MID, line=1.45, first=True)
para(tf, "That gap is the whole reason for tonight.", 20, font=SANS, color=COPPER,
     bold=True, space_before=14)
footer(s, "The Builders' Ops Studio", "01 · The Gap")

# ============================================================
# SLIDE 3 — WHY ME (photo + pillars)
# ============================================================
s = add_slide(LIGHT)
accent_bar(s)
eyebrow(s, "Why Me", 0.95, 0.95)
_, tf = textbox(s, 0.9, 1.4, 7.4, 1.8)
para(tf, "Fifteen years building", 38, font=SERIF, color=BROWN, bold=True, line=1.02, first=True)
para(tf, "other people's operations.", 38, font=SERIF, color=BROWN, bold=True, line=1.02)
_, tf = textbox(s, 0.95, 3.1, 7.0, 1.9)
para(tf, "Multimillion-dollar budgets. Teams across five continents. Then a course "
        "teaching software engineers how AI systems actually work, the real machinery "
        "behind the tools. Now that same infrastructure, built for small businesses.",
     16, font=SANS, color=BROWN_MID, line=1.4, first=True)

# pillars
pillars = [("01", "I build it", "Not slideware. Real systems that ship."),
           ("02", "Enterprise ops", "For businesses our size."),
           ("03", "I taught the AI", "No guessing on what it can do.")]
pw, ph, gap = 2.28, 1.55, 0.18
px0, py = 0.95, 5.15
for i, (k, tt, dd) in enumerate(pillars):
    px = px0 + i * (pw + gap)
    card = rect(s, px, py, pw, ph, WHITE, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
                line=SAND, line_w=1.2)
    card.adjustments[0] = 0.08
    _, tf = textbox(s, px + 0.18, py + 0.14, pw - 0.36, ph - 0.28)
    para(tf, k, 24, font=SERIF, color=CLAY, bold=True, italic=True, line=1.0, first=True)
    para(tf, tt, 15, font=SANS, color=BROWN, bold=True, space_before=4)
    para(tf, dd, 11.5, font=SANS, color=BROWN_MID, line=1.3, space_before=3)

# photo + caption bar
photo = os.path.join(ASSETS, "cynthia.jpg")
psz = 3.5
plx, pty = 9.1, 1.55
pic = s.shapes.add_picture(photo, Inches(plx), Inches(pty), Inches(psz), Inches(psz))
pic.line.color.rgb = CREAM
pic.line.width = Pt(6)
cap = rect(s, plx + 0.45, pty + psz - 0.25, psz - 0.9, 0.5, CLAY,
           shape=MSO_SHAPE.ROUNDED_RECTANGLE)
cap.adjustments[0] = 0.5
_, tf = textbox(s, plx + 0.45, pty + psz - 0.22, psz - 0.9, 0.44, anchor=MSO_ANCHOR.MIDDLE)
para(tf, "Cynthia Jones · Founder", 11, font=SANS, color=WHITE, bold=True,
     upper=True, spc=1.4, align=PP_ALIGN.CENTER, first=True)
footer(s, "The Builders' Ops Studio", "02 · Why Me")

# ============================================================
# SLIDE 4 — THE REFRAME
# ============================================================
s = add_slide(BROWN)
accent_bar(s)
eyebrow(s, "The Reframe", 0.95, 0.95)
_, tf = textbox(s, 0.9, 1.45, 11.0, 1.7)
para(tf, [("Most AI talk is about the ", {}), ("front", {"color": CLAY, "italic": True}),
          (" of your business. More followers. More leads. More customers at the door.", {})],
     33, font=SERIF, color=CREAM, italic=True, line=1.18, first=True)
_, tf = textbox(s, 0.9, 3.55, 11.0, 1.1)
para(tf, [("My work starts the moment they walk ", {}),
          ("through", {"color": CLAY, "italic": True}), (" it.", {})],
     33, font=SERIF, color=CREAM, italic=True, line=1.18, first=True)
_, tf = textbox(s, 0.95, 5.05, 10.8, 1.3)
para(tf, "“She fills the top of the funnel. I build the machine that keeps the "
        "bottom from leaking.”", 22, font=SERIF, color=CLAY, italic=True,
     line=1.3, first=True)
footer(s, "Cynthia Jones", "03 · Front vs. Back")

# ============================================================
# SLIDE 5 — BIG TRUTH
# ============================================================
s = add_slide(CREAM)
accent_bar(s)
eyebrow(s, "The Stakes", 0.95, 1.6)
_, tf = textbox(s, 0.9, 2.2, 11.4, 3.0)
para(tf, "More customers without a system underneath isn't growth.",
     54, font=SERIF, color=BROWN, bold=True, line=1.08, first=True)
para(tf, "It's just a bigger fire.", 54, font=SERIF, color=COPPER, bold=True,
     line=1.08, space_before=8)
footer(s, "The Builders' Ops Studio", "04")

# ============================================================
# SLIDE 6 — SECTION DIVIDER: SHOW
# ============================================================
s = add_slide(BROWN)
accent_bar(s)
eyebrow(s, "Show, Don't Tell", 0.95, 1.7)
_, tf = textbox(s, 0.9, 2.2, 11.4, 2.5)
para(tf, "Two real", 78, font=SERIF, color=CREAM, bold=True, line=0.98, first=True)
para(tf, "businesses.", 78, font=SERIF, color=CREAM, bold=True, line=0.98)
_, tf = textbox(s, 0.95, 5.2, 10.6, 1.0)
para(tf, "Not demos. Not mockups. Systems running right now for a Georgia bakery and "
        "a wellness coach.", 20, font=SANS, color=SAND, line=1.4, first=True)
footer(s, "Cynthia Jones", "05 · The Builds")


# ============================================================
# Build-detail slides (Baker, Coach) share a layout
# ============================================================
def beats_slide(eyebrow_txt, h1, h2, sub, beats, footer_r):
    s = add_slide(LIGHT)
    accent_bar(s)
    eyebrow(s, eyebrow_txt, 0.95, 0.85)
    _, tf = textbox(s, 0.9, 1.3, 11.4, 1.5)
    para(tf, h1, 36, font=SERIF, color=BROWN, bold=True, line=1.02, first=True)
    para(tf, h2, 36, font=SERIF, color=BROWN, bold=True, line=1.02)
    _, tf = textbox(s, 0.95, 2.95, 11.0, 0.6)
    para(tf, sub, 17, font=SANS, color=BROWN_MID, line=1.35, first=True)
    by = 3.75
    for k, strong, span in beats:
        _, tf = textbox(s, 0.95, by, 1.0, 1.0)
        para(tf, k, 30, font=SERIF, color=CLAY, bold=True, italic=True, first=True)
        _, tf = textbox(s, 1.85, by, 10.4, 1.1)
        para(tf, strong, 20, font=SANS, color=BROWN, bold=True, line=1.05, first=True)
        para(tf, span, 14.5, font=SANS, color=BROWN_MID, line=1.35, space_before=3)
        by += 1.18
    footer(s, "The Builders' Ops Studio", footer_r)
    return s


beats_slide(
    "Build One · The Baker",
    "Twenty years of demand.", "No system.",
    "Orders in her DMs. Capacity in her head. Turning away business she couldn't track.",
    [("01", "A live order form with a capacity tracker",
      "She never over-commits her oven again. The system says when she's full."),
     ("02", "Standing orders that run themselves",
      "Her regulars get a text every Thursday, the order processes automatically. "
      "Guaranteed weekly revenue on autopilot."),
     ("03", "No monthly subscription",
      "She owns it. The software and the data are hers.")],
    "06 · The Baker")

beats_slide(
    "Build Two · The Wellness Coach",
    "Real brand. Real following.", "All of it in her head.",
    "Everything ran on her phone and her memory.",
    [("01", "A client portal",
      "Her whole coaching relationship in one place: intake, sessions, progress, contracts."),
     ("02", "One-click enroll",
      "An application comes in, the right emails go out automatically. Turning an "
      "applicant into a client takes one click instead of a week."),
     ("03", "She got her hours back",
      "From running her practice in her head to running it from a dashboard.")],
    "07 · The Coach")

# ============================================================
# SLIDE 9 — WHAT THIS MEANS
# ============================================================
s = add_slide(BROWN)
accent_bar(s)
eyebrow(s, "What This Means For You", 0.95, 0.95)
_, tf = textbox(s, 0.9, 1.4, 11.4, 1.6)
para(tf, [("Built, not rented.", {})], 50, font=SERIF, color=CREAM, bold=True,
     line=1.04, first=True)
para(tf, [("And finally ", {}), ("affordable.", {"color": CLAY})],
     50, font=SERIF, color=CREAM, bold=True, line=1.04)
_, tf = textbox(s, 0.95, 3.6, 11.0, 1.2)
para(tf, "These aren't apps they rented. They were built for the specific business, "
        "the owner owns them, and they're designed to run without her touching every piece.",
     18, font=SANS, color=SAND, line=1.45, first=True)
_, tf = textbox(s, 0.95, 4.95, 11.2, 1.5)
para(tf, "Custom software used to cost fifty thousand dollars and a development team. "
        "AI is what makes building this for a small business affordable now. That's the "
        "real headline. You can finally own infrastructure that used to be only for big "
        "companies.", 18, font=SANS, color=SAND, line=1.45, first=True)
footer(s, "Cynthia Jones", "08 · Why Now")

# ============================================================
# SLIDE 10 — ONE QUESTION
# ============================================================
s = add_slide(CREAM)
accent_bar(s)
eyebrow(s, "Where To Start", 0.95, 1.25)
_, tf = textbox(s, 0.9, 1.75, 11.4, 2.8)
para(tf, "You don't need all of it at once.", 48, font=SERIF, color=BROWN, bold=True,
     line=1.08, first=True)
para(tf, "It starts with one question:", 48, font=SERIF, color=BROWN, bold=True, line=1.08)
para(tf, "where is your business breaking?", 48, font=SERIF, color=COPPER, bold=True,
     line=1.08)
_, tf = textbox(s, 0.95, 5.5, 10.8, 0.8)
para(tf, "We find that first. Then we build the one thing that fixes it.",
     20, font=SANS, color=BROWN_MID, line=1.4, first=True)
footer(s, "The Builders' Ops Studio", "09 · One Question")

# ============================================================
# SLIDE 11 — THE SERIES
# ============================================================
s = add_slide(LIGHT)
accent_bar(s)
eyebrow(s, "An Idea I'm Testing", 0.95, 1.35)
_, tf = textbox(s, 0.9, 1.85, 11.4, 1.2)
para(tf, "A small hands-on series.", 60, font=SERIF, color=BROWN, bold=True,
     line=1.02, first=True)
_, tf = textbox(s, 0.95, 3.5, 10.8, 1.0)
para(tf, "A few of us, building one real system together over a few weeks. Not a "
        "lecture. A build.", 20, font=SANS, color=BROWN_MID, line=1.4, first=True)
_, tf = textbox(s, 0.95, 4.7, 10.8, 1.0)
para(tf, "If that sounds useful, find me after and tell me. I'm listening for whether "
        "there's appetite for it.", 20, font=SANS, color=COPPER, bold=True, line=1.4,
     first=True)
footer(s, "The Builders' Ops Studio", "10 · The Series")

# ============================================================
# SLIDE 12 — THE CLOSE + QR
# ============================================================
s = add_slide(BROWN)
accent_bar(s)
eyebrow(s, "The Ask", 0.95, 1.1)
_, tf = textbox(s, 0.9, 1.55, 6.9, 2.6)
para(tf, "If your business is", 40, font=SERIF, color=CREAM, bold=True, line=1.05, first=True)
para(tf, "running you,", 40, font=SERIF, color=CREAM, bold=True, line=1.05)
para(tf, "let's fix that.", 40, font=SERIF, color=CREAM, bold=True, line=1.05)
_, tf = textbox(s, 0.95, 4.35, 6.6, 1.3)
para(tf, "A free 45-minute discovery call. We sit down, I help you see where the gaps "
        "are. No pitch.", 18, font=SANS, color=SAND, line=1.45, first=True)
_, tf = textbox(s, 0.95, 5.75, 6.6, 1.0)
para(tf, "The code books it. Grab a slot tonight.", 16, font=SANS, color=SAND,
     line=1.5, first=True)
para(tf, "cynthia@thebuildersopsstudio.com", 16, font=SANS, color=CLAY, bold=True,
     space_before=3)

# QR card
qr = os.path.join(ASSETS, "qr-booking.png")
cw, ch = 4.0, 4.5
cx, cy = 8.4, 1.4
card = rect(s, cx, cy, cw, ch, CREAM, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
card.adjustments[0] = 0.06
qsz = 3.0
s.shapes.add_picture(qr, Inches(cx + (cw - qsz) / 2), Inches(cy + 0.35),
                     Inches(qsz), Inches(qsz))
_, tf = textbox(s, cx, cy + 3.55, cw, 0.4)
para(tf, "Scan to book your call", 13, font=SANS, color=COPPER, bold=True, upper=True,
     spc=1.8, align=PP_ALIGN.CENTER, first=True)
_, tf = textbox(s, cx, cy + 3.95, cw, 0.35)
para(tf, "thebuildersopsstudio.com", 12, font=SANS, color=BROWN_MID,
     align=PP_ALIGN.CENTER, first=True)
footer(s, "Cynthia Jones", "11 · Book A Call")

# ============================================================
# SLIDE 13 — THANK YOU
# ============================================================
s = add_slide(BROWN)
accent_bar(s)
wordmark(s, 0.95, 0.95, CREAM, scale=1.0)
_, tf = textbox(s, 0.95, 2.75, 6, 0.4)
para(tf, "Strategy · Systems · Structure", 12, font=SANS, color=CLAY, bold=True,
     upper=True, spc=2.2, first=True)
_, tf = textbox(s, 0.9, 3.4, 11.4, 2.0)
para(tf, "I build the systems", 60, font=SERIF, color=CREAM, bold=True, line=1.0, first=True)
para(tf, "your business runs on.", 60, font=SERIF, color=CREAM, bold=True, line=1.0)
_, tf = textbox(s, 0.95, 5.95, 11.0, 0.7)
para(tf, "Cynthia Jones · The Builders' Ops Studio. Thank you.",
     20, font=SANS, color=SAND, line=1.4, first=True)
footer(s, "cynthia@thebuildersopsstudio.com", "thebuildersopsstudio.com")

out = os.path.join(HERE, "Beyond-Prompting.pptx")
prs.save(out)
print("Saved:", out, "| slides:", len(prs.slides._sldIdLst))
